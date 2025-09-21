from rest_framework import generics, status, permissions
from rest_framework.decorators import api_view, permission_classes
from rest_framework.response import Response
from django.shortcuts import get_object_or_404
from django.db.models import Count, Q
from django.utils import timezone
from django.http import HttpResponse
from django.template.loader import render_to_string
from django.contrib.auth.models import User
from apps.team.models import Team, TeamMembership
from apps.actions.models import ActionDb
from apps.education.models import EducationDb
from apps.notifications.models import Notification
from datetime import timedelta
from django.db import transaction
from weasyprint import HTML, CSS
from channels.layers import get_channel_layer
from asgiref.sync import async_to_sync
from .models import Form, FormContent, FormEditHistory, FormEditSession
from .serializers import (
    FormListSerializer,
    FormDetailSerializer, 
    CreateFormSerializer,
    FormContentSerializer,
    FormEditHistorySerializer,
    FormEditSessionSerializer,
    CreatePersonalFormSerializer
)
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import io
import json
import markdown

class TeamFormListCreateView(generics.ListCreateAPIView):
    """Team form list and creation API"""
    permission_classes = [permissions.IsAuthenticated]
    
    def get_queryset(self):
        """Obtain the list of forms for the specified team"""
        team_id = self.kwargs.get('team_id')
        team = get_object_or_404(Team, id=team_id)
        
        # Verify whether the user is a team member
        membership = TeamMembership.objects.filter(
            user=self.request.user,
            team=team
        ).first()
        
        if not membership:
            return Form.objects.none()
        
        return Form.objects.filter(team=team).select_related(
            'created_by', 'last_modified_by', 'team'
        )
    
    def get_serializer_class(self):
        if self.request.method == 'POST':
            return CreateFormSerializer
        return FormListSerializer
    
    def get_serializer_context(self):
        context = super().get_serializer_context()
        context['team_id'] = self.kwargs.get('team_id')
        return context
    
    def create(self, request, *args, **kwargs):
        """Create form - Requires edit or owner privileges"""
        team_id = kwargs.get('team_id')
        team = get_object_or_404(Team, id=team_id)
        
        # Check permissions
        membership = TeamMembership.objects.filter(
            user=request.user,
            team=team,
            role__in=['owner', 'edit']
        ).first()
        
        if not membership:
            return Response(
                {'error': 'Only team owner or edit member can create forms'}, 
                status=status.HTTP_403_FORBIDDEN
            )
        
        return super().create(request, *args, **kwargs)

class FormDetailView(generics.RetrieveUpdateDestroyAPIView):
    """Form Details API"""
    permission_classes = [permissions.IsAuthenticated]
    serializer_class = FormDetailSerializer
    
    def get_queryset(self):
        """Obtain the forms that the user has permission to access"""
        return Form.objects.select_related(
            'created_by', 'last_modified_by', 'team'
        )
    
    def get_object(self):
        """Obtain the form object and check the permissions"""
        obj = super().get_object()

        if obj.team is None:
            if obj.created_by != self.request.user:
                from django.http import Http404
                raise Http404("Form not found")
            return obj
        
        # Check whether the user is a team member
        membership = TeamMembership.objects.filter(
            user=self.request.user,
            team=obj.team
        ).first()
        
        if not membership:
            from django.http import Http404
            raise Http404("Form not found")
        
        return obj
    
    def update(self, request, *args, **kwargs):
        """Update form - Requires write or admin privileges"""
        form = self.get_object()

        if form.team is None:
            if form.created_by != request.user:
                return Response(
                    {'error': 'Only the form creator can update this personal form'}, 
                    status=status.HTTP_403_FORBIDDEN
                )
        else:
            permission = form.get_user_permission(request.user)
            if permission not in ['write', 'admin']:
                return Response(
                    {'error': 'Insufficient permissions to update this form'}, 
                    status=status.HTTP_403_FORBIDDEN
                )
        
        # Update last modifier
        if hasattr(form, 'last_modified_by'):
            form.last_modified_by = request.user
            form.save(update_fields=['last_modified_by'])
        
        return super().update(request, *args, **kwargs)
    
    def destroy(self, request, *args, **kwargs):
        """Delete form - Requires admin privileges"""
        form = self.get_object()

        if form.team is None:
            if form.created_by != request.user:
                return Response(
                    {'error': 'Only the form creator can delete this personal form'}, 
                    status=status.HTTP_403_FORBIDDEN
                )
        else:
            permission = form.get_user_permission(request.user)
            if permission != 'admin':
                return Response(
                    {'error': 'Only form admin can delete this form'}, 
                    status=status.HTTP_403_FORBIDDEN
                )
        
        return super().destroy(request, *args, **kwargs)

@api_view(['GET'])
@permission_classes([permissions.IsAuthenticated])
def team_form_stats(request, team_id):
    """Obtain the statistics of the team form"""
    team = get_object_or_404(Team, id=team_id)
    
    # Verify permissions
    membership = TeamMembership.objects.filter(
        user=request.user,
        team=team
    ).first()
    
    if not membership:
        return Response(
            {'error': 'You are not a member of this team'}, 
            status=status.HTTP_403_FORBIDDEN
        )
    
    # Obtain statistical data
    forms = Form.objects.filter(team=team)
    
    stats = {
        'total_forms': forms.count(),
        'active_forms': forms.filter(status='active').count(),
        'locked_forms': forms.filter(status='locked').count(),
        'archived_forms': forms.filter(status='archived').count(),
        'forms_by_type': {
            'action': forms.filter(type='action').count(),
            'education': forms.filter(type='education').count(),
            'blank': forms.filter(type='blank').count(),
            'ida': forms.filter(type='ida').count(),
        }
    }
    
    return Response(stats)

@api_view(['PATCH'])
@permission_classes([permissions.IsAuthenticated])
def toggle_form_lock(request, form_id):
    """Switch the form lock status"""
    form = get_object_or_404(Form, id=form_id)
    
    permission = form.get_user_permission(request.user)
    if permission != 'admin':
        return Response(
            {'error': 'Only form admin can lock/unlock forms'}, 
            status=status.HTTP_403_FORBIDDEN
        )
    
    # Switching status
    if form.status == 'locked':
        form.status = 'active'
        message = 'Form unlocked successfully'
    elif form.status == 'active':
        form.status = 'locked'
        message = 'Form locked successfully'
    else:
        return Response(
            {'error': 'Can only lock/unlock active or locked forms'}, 
            status=status.HTTP_400_BAD_REQUEST
        )
    
    form.last_modified_by = request.user
    form.save(update_fields=['status', 'last_modified_by'])
    
    serializer = FormDetailSerializer(form, context={'request': request})
    return Response({
        'message': message,
        'form': serializer.data
    })

@api_view(['POST'])
@permission_classes([permissions.IsAuthenticated])
def duplicate_form(request, form_id):
    """Copy the form"""
    original_form = get_object_or_404(Form, id=form_id)
    
    # Check permissions
    permission = original_form.get_user_permission(request.user)
    if not permission:
        return Response(
            {'error': 'You do not have permission to access this form'}, 
            status=status.HTTP_403_FORBIDDEN
        )
    
    # Check if user has permission to create in the team
    membership = TeamMembership.objects.filter(
        user=request.user,
        team=original_form.team,
        role__in=['owner', 'edit']
    ).first()
    
    if not membership:
        return Response(
            {'error': 'Only team owner or edit member can duplicate forms'}, 
            status=status.HTTP_403_FORBIDDEN
        )
    
    # Create a duplicate copy
    duplicated_form = Form.objects.create(
        title=f"{original_form.title} (Copy)",
        description=original_form.description,
        type=original_form.type,
        status='active',
        team=original_form.team,
        created_by=request.user,
        last_modified_by=request.user,
        allow_anonymous=original_form.allow_anonymous,
        allow_multiple_submissions=original_form.allow_multiple_submissions,
        require_login=original_form.require_login,
        is_public=False,
        is_template=False,
        deadline=original_form.deadline,
        max_responses=original_form.max_responses,
    )
    
    serializer = FormDetailSerializer(duplicated_form, context={'request': request})
    return Response({
        'message': 'Form duplicated successfully',
        'form': serializer.data
    }, status=status.HTTP_201_CREATED)

class CollaborativeFormDetailView(generics.RetrieveUpdateAPIView):
    """Collaborative form detail and update API"""
    permission_classes = [permissions.IsAuthenticated]
    serializer_class = FormContentSerializer
    
    def get_object(self):
        form_id = self.kwargs.get('form_id')
        form = get_object_or_404(Form, id=form_id)
        
        # Check permissions
        if not self.has_form_access(form):
            from django.http import Http404
            raise Http404("Form not found")
        
        content, created = FormContent.objects.get_or_create(
            form=form,
            defaults={
                'version': 1,
                'title': form.title,
            }
        )

        content.form_status = form.status
        return content
    
    def has_form_access(self, form):
        """Check if user has access to the form"""
        if form.team is None:
            return form.created_by == self.request.user

        from apps.team.models import TeamMembership
        membership = TeamMembership.objects.filter(
            user=self.request.user,
            team=form.team
        ).first()
        return membership is not None
    
    def get_editable_fields(self, form_type):
        """Return editable fields based on form type"""
        common_fields = ['title', 'description', 'location', 'organization', 'year', 'sdgs_related', 'source', 'link']
        
        if form_type == 'education':
            return common_fields + [
                'aims', 'learning_outcomes', 'type_label', 
                'related_discipline', 'useful_industries'
            ]
        elif form_type == 'action':
            return common_fields + [
                'actions', 'action_detail', 'level', 'individual_organization',
                'related_industry', 'digital_actions',
                'source_descriptions', 'award', 'source_links', 
                'additional_notes', 'award_descriptions'
            ]
        elif form_type == 'blank':
            return ['title', 'description', 'free_content']
        elif form_type == 'ida':
            return [
                'title', 'description', 'designer_names', 
                'current_role_affiliation', 'impact_project_name', 'main_challenge',
                'project_description', 'selected_sdgs', 'impact_types', 
                'project_importance', 'existing_example',
                'implementation_step1', 'implementation_step2', 'implementation_step3',
                'implementation_step4', 'implementation_step5', 'implementation_step6',
                'resources_partnerships', 'skills_capabilities', 'impact_avenues',
                'risks_inhibitors', 'mitigation_strategies'
            ]
        return common_fields
    
    def update(self, request, *args, **kwargs):
        """Update form content"""
        content = self.get_object()
        field_name = request.data.get('field_name')
        field_value = request.data.get('field_value')
        
        if not field_name:
            return Response(
                {'error': 'field_name is required'},
                status=status.HTTP_400_BAD_REQUEST
            )
        
        # Check if the field is editable for this form type
        editable_fields = self.get_editable_fields(content.form.type)
        if field_name not in editable_fields:
            return Response(
                {'error': f'Field {field_name} is not editable for this form type'},
                status=status.HTTP_400_BAD_REQUEST
            )
        
        # Check if user has permission to edit this field
        if not self.has_field_edit_permission(content.form, field_name):
            return Response(
                {'error': 'No permission to edit this field'},
                status=status.HTTP_403_FORBIDDEN
            )
        
        # Log edit history
        old_value = getattr(content, field_name, '')
        
        with transaction.atomic():
            setattr(content, field_name, field_value)
            content.version += 1
            content.save()

            # If updating title, also sync it to Form model

            if field_name == 'title' and field_value:
                content.form.title = field_value
                content.form.save(update_fields=['title'])
            
            # Log edit history
            FormEditHistory.objects.create(
                form=content.form,
                user=request.user,
                field_name=field_name,
                old_value=old_value,
                new_value=field_value,
                version=content.version
            )
        
        self.broadcast_update(content.form, field_name, field_value, request.user)
        
        return Response({
            'success': True,
            'version': content.version,
            'field_name': field_name,
            'field_value': field_value
        })
    
    def has_field_edit_permission(self, form, field_name):
        """Check if user has permission to edit this field"""
        if form.team is None:
            return form.created_by == self.request.user
        
        membership = TeamMembership.objects.filter(
            user=self.request.user,
            team=form.team
        ).first()
        
        return membership and membership.role in ['owner', 'edit']
    
    def broadcast_update(self, form, field_name, field_value, user):
        """Broadcast update to other users via WebSocket"""
        channel_layer = get_channel_layer()
        async_to_sync(channel_layer.group_send)(
            f"form_{form.id}",
            {
                'type': 'field_update',
                'field_name': field_name,
                'field_value': field_value,
                'user_id': user.id,
                'user_name': user.username,
                'timestamp': json.dumps(timezone.now(), default=str)
            }
        )

def normalize_sdg_data(sdg_value):
    """Normalize SDG data to a comma-separated string"""
    if not sdg_value:
        return ""
    
    if isinstance(sdg_value, str) and sdg_value.startswith('['):
        try:
            import json
            sdg_list = json.loads(sdg_value)
            if isinstance(sdg_list, list):
                return ",".join([str(sdg).strip() for sdg in sdg_list])
        except:
            pass
    
    elif isinstance(sdg_value, list):
        return ",".join([str(sdg).strip() for sdg in sdg_value])

@api_view(['PATCH'])
@permission_classes([permissions.IsAuthenticated])
def collaborative_form_batch_update(request, form_id):
    """Batch update multiple fields of the form"""
    form = get_object_or_404(Form, id=form_id)
    
    if form.team is None:
        if form.created_by != request.user:
            return Response(
                {'error': 'Only the form creator can edit this personal form'}, 
                status=status.HTTP_403_FORBIDDEN
            )
    else:
        # Check permissions
        membership = TeamMembership.objects.filter(
            user=request.user,
            team=form.team
        ).first()
        
        if not membership:
            return Response(
                {'error': 'You are not a member of this team'}, 
                status=status.HTTP_403_FORBIDDEN
            )
        
        if membership.role not in ['owner', 'edit']:
            return Response(
                {'error': 'No permission to edit this form'},
                status=status.HTTP_403_FORBIDDEN
            )
        

    content, created = FormContent.objects.get_or_create(
        form=form,
        defaults={
            'version': 1,
            'title': form.title,
        }
    )
        
    changes = request.data.get('changes', {})
    if not changes:
        return Response(
            {'error': 'No changes provided'},
            status=status.HTTP_400_BAD_REQUEST
        )
    
    def get_field_default_value(field_name, current_value):
        """Return appropriate default value based on field type"""
        if current_value is None:
            boolean_fields = ['digital_actions', 'award']
            integer_fields = ['level', 'individual_organization', 'year']
            
            if field_name in boolean_fields:
                return False
            elif field_name in integer_fields:
                return 0
            else:
                return ''
        return current_value
    
    # Batch update multiple fields of the form
    with transaction.atomic():
        old_values = {}
        
        for field_name, field_value in changes.items():
            if hasattr(content, field_name):
                current_value = getattr(content, field_name, None)

                old_values[field_name] = get_field_default_value(field_name, current_value)
                

                if field_name in ['level', 'individual_organization']:

                    if field_value == '' or field_value is None:
                        field_value = None
                    else:
                        try:
                            field_value = int(field_value)
                        except (ValueError, TypeError):
                            field_value = None
                elif field_name in ['digital_actions', 'award']:
                    if isinstance(field_value, str):
                        if field_value == '1' or field_value.lower() == 'true':
                            field_value = True
                        elif field_value == '0' or field_value.lower() == 'false':
                            field_value = False
                        else:
                            field_value = bool(field_value)
                elif field_name in ['sdgs_related', 'selected_sdgs']:
                    field_value = normalize_sdg_data(field_value)
                
                setattr(content, field_name, field_value)
                

                if field_name == 'title' and field_value:
                    content.form.title = field_value
                    content.form.save(update_fields=['title'])
        

        content.version += 1
        content.save()
        

        history_records = []
        for field_name, field_value in changes.items():
            if field_name in old_values:
                old_val = old_values[field_name]
                new_val = field_value
                

                if old_val is None:
                    old_val = ''
                if new_val is None:
                    new_val = ''
                
           
                history_records.append(FormEditHistory(
                    form=content.form,
                    user=request.user,
                    field_name=field_name,
                    old_value=str(old_val),
                    new_value=str(new_val),
                    version=content.version
                ))
        
        if history_records:
            FormEditHistory.objects.bulk_create(history_records)
    
    # Broadcast batch update via WebSocket
    channel_layer = get_channel_layer()
    async_to_sync(channel_layer.group_send)(
        f"form_{form.id}",
        {
            'type': 'batch_update',
            'changes': changes,
            'version': content.version,
            'user_id': request.user.id,
            'user_name': request.user.username,
            'timestamp': json.dumps(timezone.now(), default=str)
        }
    )
    
    return Response({
        'success': True,
        'version': content.version,
        'changes_count': len(changes),
        'updated_fields': list(changes.keys())
    })

@api_view(['POST'])
@permission_classes([permissions.IsAuthenticated])
def start_edit_session(request, form_id):
    """Start an edit session"""
    form = get_object_or_404(Form, id=form_id)
    field_name = request.data.get('field_name')
    
    if form.team is None:
        if form.created_by != request.user:
            return Response(
                {'error': 'Only the form creator can edit this personal form'}, 
                status=status.HTTP_403_FORBIDDEN
            )
        
    session, created = FormEditSession.objects.update_or_create(
        form=form,
        user=request.user,
        defaults={
            'field_name': field_name,
            'is_active': True,
            'cursor_position': request.data.get('cursor_position', 0)
        }
    )
    
    # Broadcast editing status

    channel_layer = get_channel_layer()
    async_to_sync(channel_layer.group_send)(
        f"form_{form.id}",
        {
            'type': 'user_editing',
            'user_id': request.user.id,
            'user_name': request.user.username,
            'field_name': field_name,
            'cursor_position': session.cursor_position
        }
    )
    
    return Response({'success': True})

@api_view(['POST'])
@permission_classes([permissions.IsAuthenticated])
def end_edit_session(request, form_id):
    """End edit session"""
    form = get_object_or_404(Form, id=form_id)

    if form.team is None:
        if form.created_by != request.user:
            return Response(
                {'error': 'Only the form creator can edit this personal form'}, 
                status=status.HTTP_403_FORBIDDEN
            )
    
    FormEditSession.objects.filter(
        form=form,
        user=request.user
    ).update(is_active=False)
    
    # Broadcast that editing has ended

    channel_layer = get_channel_layer()
    async_to_sync(channel_layer.group_send)(
        f"form_{form.id}",
        {
            'type': 'user_stopped_editing',
            'user_id': request.user.id,
            'user_name': request.user.username
        }
    )
    
    return Response({'success': True})

@api_view(['GET'])
@permission_classes([permissions.IsAuthenticated])
def get_active_editors(request, form_id):
    """Get list of currently active editors"""
    form = get_object_or_404(Form, id=form_id)

    if form.team is None:
        if form.created_by != request.user:
            return Response(
                {'error': 'Only the form creator can access this personal form'}, 
                status=status.HTTP_403_FORBIDDEN
            )
    
    sessions = FormEditSession.objects.filter(
        form=form,
        is_active=True
    ).select_related('user')
    
    serializer = FormEditSessionSerializer(sessions, many=True)
    return Response(serializer.data)

@api_view(['GET'])
@permission_classes([permissions.IsAuthenticated])
def export_form_pdf(request, form_id):
    form = get_object_or_404(Form, id=form_id)
    content = get_object_or_404(FormContent, form=form)
    
    if form.type == 'blank':
        template_name = 'exports/blank_form.html'
        context = {
            'title': content.title or form.title,
            'description': content.description or form.description,
            'content_html': markdown.markdown(content.free_content) if content.free_content else '',
            'form': form,
        }
    elif form.type == 'action':
        template_name = 'exports/action_form.html'
        context = {
            'form': form,
            'content': content,
            'sections': [
                ('Basic Information', [
                    ('Title', content.title),
                    ('Description', content.description),
                    ('Location', content.location),
                    ('Organization', content.organization),
                    ('Year', content.year),
                ]),
                ('Action Details', [
                    ('Actions', content.actions),
                    ('Action Detail', content.action_detail),
                    ('Level', content.level),
                    ('Individual/Organization', content.individual_organization),
                    ('Related Industry', content.related_industry),
                    ('Digital Actions', 'Yes' if content.digital_actions else 'No'),
                ]),
                ('Additional Information', [
                    ('SDGs Related', content.sdgs_related),
                    ('Source Descriptions', content.source_descriptions),
                    ('Award', 'Yes' if content.award else 'No'),
                    ('Source Links', content.source_links),
                    ('Additional Notes', content.additional_notes),
                    ('Award Descriptions', content.award_descriptions),
                ])
            ]
        }
    elif form.type == 'ida':
        template_name = 'exports/ida_form.html'
        

        selected_sdgs_display = []
        if content.selected_sdgs:
            try:
                sdg_list = json.loads(content.selected_sdgs) if isinstance(content.selected_sdgs, str) else content.selected_sdgs
                for sdg_num in sdg_list:
                    selected_sdgs_display.append(f"SDG {sdg_num}")
            except:
                selected_sdgs_display = [content.selected_sdgs] if content.selected_sdgs else []
        
        impact_types_display = []
        if content.impact_types:
            try:
                impact_dict = json.loads(content.impact_types) if isinstance(content.impact_types, str) else content.impact_types
                impact_types_display = [f"{k}: {v}" for k, v in impact_dict.items()]
            except:
                impact_types_display = [content.impact_types] if content.impact_types else []
        
        context = {
            'form': form,
            'content': content,
            'selected_sdgs_display': ', '.join(selected_sdgs_display),
            'impact_types_display': ', '.join(impact_types_display),
            'sections': [
                ('Basic Information', [
                    ('Designer Names', content.designer_names),
                    ('Current Role and Affiliation', content.current_role_affiliation),
                ]),
                ('Project Overview', [
                    ('Impact Project Name', content.impact_project_name),
                    ('Main Challenge', content.main_challenge),
                    ('Project Description', content.project_description),
                    ('Selected SDGs', ', '.join(selected_sdgs_display)),
                    ('Impact Types', ', '.join(impact_types_display)),
                ]),
                ('Project Analysis', [
                    ('Project Importance', content.project_importance),
                    ('Existing Example', content.existing_example),
                ]),
                ('Implementation Steps', [
                    ('Step 1', content.implementation_step1),
                    ('Step 2', content.implementation_step2),
                    ('Step 3', content.implementation_step3),
                    ('Step 4', content.implementation_step4),
                    ('Step 5', content.implementation_step5),
                    ('Step 6', content.implementation_step6),
                ]),
                ('Resources and Assessment', [
                    ('Resources & Partnerships', content.resources_partnerships),
                    ('Skills & Capabilities', content.skills_capabilities),
                    ('Impact Avenues', content.impact_avenues),
                    ('Risks & Inhibitors', content.risks_inhibitors),
                    ('Mitigation Strategies', content.mitigation_strategies),
                ])
            ]
        }
    else:  # education
        template_name = 'exports/education_form.html'
        context = {
            'form': form,
            'content': content,
            'sections': [
                ('Basic Information', [
                    ('Title', content.title),
                    ('Description', content.description),
                    ('Location', content.location),
                    ('Organization', content.organization),
                    ('Year', content.year),
                ]),
                ('Educational Content', [
                    ('Aims', content.aims),
                    ('Learning Outcomes', content.learning_outcomes),
                    ('Type Label', content.type_label),
                    ('Related Discipline', content.related_discipline),
                    ('Useful Industries', content.useful_industries),
                ]),
                ('Resource Information', [
                    ('SDGs Related', content.sdgs_related),
                    ('Source', content.source),
                    ('Link', content.link),
                ])
            ]
        }
    
    # HTML
    html_string = render_to_string(template_name, context)
    
    # PDF
    pdf = HTML(string=html_string, base_url=request.build_absolute_uri()).write_pdf()
    
    response = HttpResponse(pdf, content_type='application/pdf')
    filename = f"{form.title or 'form'}_{form.id}.pdf"
    response['Content-Disposition'] = f'attachment; filename="{filename}"'
    return response

@api_view(['GET'])
@permission_classes([permissions.IsAuthenticated])
def export_ida_ppt(request, form_id):
    """Export IDA form to PPT format"""
    form = get_object_or_404(Form, id=form_id)
    
    if form.type != 'ida':
        return Response(
            {'error': 'PPT export is only available for IDA forms'}, 
            status=status.HTTP_400_BAD_REQUEST
        )
    
    content = get_object_or_404(FormContent, form=form)
    
    # PPT
    ppt_file = create_ida_presentation(form, content)
    
    # return PPT
    response = HttpResponse(
        ppt_file.getvalue(),
        content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation'
    )
    safe_title = (content.title or form.title or 'IDA_Form').replace(' ', '_')[:50]
    filename = f"{safe_title}_{form.id}.pptx"
    response['Content-Disposition'] = f'attachment; filename="{filename}"'
    return response

def create_ida_presentation(form, content):
    """Create PowerPoint presentation for IDA analysis"""
    prs = Presentation()
    
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    create_title_slide(prs, form, content)
    
    create_overview_slide(prs, content)
    
    create_roadmap_slide(prs, content)
    
    create_assessment_slide(prs, content)
    
    create_risk_slide(prs, content)
    
    ppt_file = io.BytesIO()
    prs.save(ppt_file)
    ppt_file.seek(0)
    
    return ppt_file

def create_title_slide(prs, form, content):
    """Create title slide"""
    slide_layout = prs.slide_layouts[6]  
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(252, 251, 247) 
    
    header_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.33), Inches(0.8)
    )
    header_fill = header_bar.fill
    header_fill.solid()
    header_fill.fore_color.rgb = RGBColor(139, 134, 128) 
    

    title_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.33), Inches(1.2))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = "IMPACT DESIGN ANALYSIS"
    title_p.font.size = Pt(42)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(87, 83, 78) 
    title_p.alignment = PP_ALIGN.CENTER
    

    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11.33), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.text = "SDG Action Plan & Impact Assessment"
    subtitle_p.font.size = Pt(20)
    subtitle_p.font.color.rgb = RGBColor(120, 113, 108)  
    subtitle_p.alignment = PP_ALIGN.CENTER
    

    info_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2.5), Inches(4), Inches(8.33), Inches(2.5)
    )
    info_fill = info_card.fill
    info_fill.solid()
    info_fill.fore_color.rgb = RGBColor(255, 254, 251)  
    info_card.line.color.rgb = RGBColor(214, 211, 209)  
    info_card.line.width = Pt(1)
    

    info_frame = info_card.text_frame
    info_frame.margin_left = Inches(0.5)
    info_frame.margin_top = Inches(0.3)
    info_frame.margin_right = Inches(0.5)
    info_frame.margin_bottom = Inches(0.3)
    

    p1 = info_frame.paragraphs[0]
    p1.text = content.title or "IDA Analysis Document"
    p1.font.size = Pt(18)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(87, 83, 78)
    p1.alignment = PP_ALIGN.CENTER
    

    if content.description:
        p2 = info_frame.add_paragraph()
        desc_text = content.description[:80] + "..." if len(content.description) > 80 else content.description
        p2.text = desc_text
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(120, 113, 108)
        p2.alignment = PP_ALIGN.CENTER
    

    p3 = info_frame.add_paragraph()
    p3.text = f"Designer: {content.designer_names or 'Not specified'}"
    p3.font.size = Pt(12)
    p3.font.color.rgb = RGBColor(146, 142, 138)
    p3.alignment = PP_ALIGN.CENTER
    

    p4 = info_frame.add_paragraph()
    p4.text = f"Created: {form.created_at.strftime('%B %d, %Y')}"
    p4.font.size = Pt(12)
    p4.font.color.rgb = RGBColor(146, 142, 138)
    p4.alignment = PP_ALIGN.CENTER
    

    if content.selected_sdgs:
        try:
            sdg_list = json.loads(content.selected_sdgs) if isinstance(content.selected_sdgs, str) else content.selected_sdgs
            if sdg_list:
                sdg_text = f"SDGs: {', '.join([f'SDG {sdg}' for sdg in sdg_list[:3]])}"
                p5 = info_frame.add_paragraph()
                p5.text = sdg_text
                p5.font.size = Pt(12)
                p5.font.bold = True
                p5.font.color.rgb = RGBColor(139, 134, 128)
                p5.alignment = PP_ALIGN.CENTER
        except:
            pass

def create_overview_slide(prs, content):
    """Create overview slide (card layout)"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(252, 251, 247)
    

    header_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.33), Inches(0.6)
    )
    header_fill = header_bar.fill
    header_fill.solid()
    header_fill.fore_color.rgb = RGBColor(139, 134, 128)
    

    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(11.33), Inches(0.8))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = "PROJECT OVERVIEW"
    title_p.font.size = Pt(28)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(87, 83, 78)
    title_p.alignment = PP_ALIGN.CENTER
    

    cards = [
        {
            'title': 'Main Challenge',
            'content': content.main_challenge or 'Not specified',
            'color': RGBColor(169, 156, 141),  
            'pos': (1, 2.2, 3.7, 2)
        },
        {
            'title': 'Project Description', 
            'content': content.project_description or 'Not specified',
            'color': RGBColor(147, 149, 151),  
            'pos': (4.9, 2.2, 3.7, 2)
        },
        {
            'title': 'Why Important',
            'content': content.project_importance or 'Not specified', 
            'color': RGBColor(156, 153, 136),  
            'pos': (8.8, 2.2, 3.7, 2)
        },
        {
            'title': 'Existing Example',
            'content': content.existing_example or 'Not specified',
            'color': RGBColor(176, 151, 128),  
            'pos': (1, 4.5, 3.7, 2)
        },
        {
            'title': 'Resources Needed',
            'content': content.resources_partnerships or 'Not specified',
            'color': RGBColor(163, 141, 156),  
            'pos': (4.9, 4.5, 3.7, 2)
        },
        {
            'title': 'Impact Pathways',
            'content': content.impact_avenues or 'Not specified',
            'color': RGBColor(153, 156, 141),  
            'pos': (8.8, 4.5, 3.7, 2)
        }
    ]
    

    for card in cards:
        card_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(card['pos'][0]), Inches(card['pos'][1]),
            Inches(card['pos'][2]), Inches(card['pos'][3])
        )
        

        fill = card_shape.fill
        fill.solid()
        fill.fore_color.rgb = card['color']
        

        card_shape.line.color.rgb = RGBColor(255, 254, 251)
        card_shape.line.width = Pt(2)
        
 
        text_frame = card_shape.text_frame
        text_frame.margin_left = Inches(0.2)
        text_frame.margin_right = Inches(0.2)
        text_frame.margin_top = Inches(0.2)
        text_frame.margin_bottom = Inches(0.2)
        

        title_p = text_frame.paragraphs[0]
        title_p.text = card['title']
        title_p.font.size = Pt(14)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(255, 254, 251)
        title_p.alignment = PP_ALIGN.CENTER
        

        content_p = text_frame.add_paragraph()
        content_text = card['content'][:80] + "..." if len(card['content']) > 80 else card['content']
        content_p.text = content_text
        content_p.font.size = Pt(10)
        content_p.font.color.rgb = RGBColor(255, 254, 251)

def create_roadmap_slide(prs, content):
    """Create implementation roadmap slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(252, 251, 247)
    

    header_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.33), Inches(0.6)
    )
    header_fill = header_bar.fill
    header_fill.solid()
    header_fill.fore_color.rgb = RGBColor(139, 134, 128)
    

    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(11.33), Inches(0.8))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = "IMPLEMENTATION ROADMAP"
    title_p.font.size = Pt(28)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(87, 83, 78)
    title_p.alignment = PP_ALIGN.CENTER
    

    steps = [

        {'title': 'Step 1', 'content': content.implementation_step1, 'pos': (2.2, 2.2), 'color': RGBColor(169, 156, 141)},
        {'title': 'Step 2', 'content': content.implementation_step2, 'pos': (5.2, 2.2), 'color': RGBColor(147, 149, 151)}, 
        {'title': 'Step 3', 'content': content.implementation_step3, 'pos': (8.2, 2.2), 'color': RGBColor(156, 153, 136)},

        {'title': 'Step 4', 'content': content.implementation_step4, 'pos': (8.2, 4.8), 'color': RGBColor(176, 151, 128)},
        {'title': 'Step 5', 'content': content.implementation_step5, 'pos': (5.2, 4.8), 'color': RGBColor(163, 141, 156)},
        {'title': 'Step 6', 'content': content.implementation_step6, 'pos': (2.2, 4.8), 'color': RGBColor(153, 156, 141)}
    ]
    

    line1 = slide.shapes.add_connector(
        1,
        Inches(4.0), Inches(2.9),  
        Inches(5.2), Inches(2.9)  
    )
    line1.line.color.rgb = RGBColor(193, 181, 171)
    line1.line.width = Pt(3)
    

    line2 = slide.shapes.add_connector(
        1,
        Inches(7.0), Inches(2.9),  
        Inches(8.2), Inches(2.9)   
    )
    line2.line.color.rgb = RGBColor(193, 181, 171)
    line2.line.width = Pt(3)
    

    vertical_line = slide.shapes.add_connector(
        1,
        Inches(9.1), Inches(4.0), 
        Inches(9.1), Inches(4.8)  
    )
    vertical_line.line.color.rgb = RGBColor(193, 181, 171)
    vertical_line.line.width = Pt(3)
    

    line4 = slide.shapes.add_connector(
        1,
        Inches(8.2), Inches(5.7),  
        Inches(7.0), Inches(5.7)   
    )
    line4.line.color.rgb = RGBColor(193, 181, 171)
    line4.line.width = Pt(3)
    

    line5 = slide.shapes.add_connector(
        1,
        Inches(5.2), Inches(5.7),  
        Inches(4.0), Inches(5.7)  
    )
    line5.line.color.rgb = RGBColor(193, 181, 171)
    line5.line.width = Pt(3)
    

    for step in steps:

        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(step['pos'][0]), Inches(step['pos'][1]),
            Inches(1.8), Inches(1.8)
        )
        
        fill = circle.fill
        fill.solid()
        fill.fore_color.rgb = step['color']
        

        circle.line.color.rgb = RGBColor(255, 254, 251)
        circle.line.width = Pt(4)
        

        text_frame = circle.text_frame
        text_frame.clear()  
        text_frame.margin_left = Inches(0.05)
        text_frame.margin_right = Inches(0.05) 
        text_frame.margin_top = Inches(0.15)
        text_frame.margin_bottom = Inches(0.05)
        text_frame.word_wrap = True
        

        title_p = text_frame.paragraphs[0]
        title_p.text = step['title']
        title_p.font.size = Pt(14)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(255, 254, 251)
        title_p.alignment = PP_ALIGN.CENTER
        

        content_p = text_frame.add_paragraph()
        step_content = step['content'] or "Not specified"
        

        if len(step_content) > 35:
            words = step_content.split()
            truncated = []
            char_count = 0
            for word in words:
                if char_count + len(word) + 1 <= 35:  # +1 for space
                    truncated.append(word)
                    char_count += len(word) + 1
                else:
                    break
            preview_text = ' '.join(truncated) + "..."
        else:
            preview_text = step_content
            
        content_p.text = preview_text
        content_p.font.size = Pt(10)
        content_p.font.color.rgb = RGBColor(255, 254, 251)
        content_p.alignment = PP_ALIGN.CENTER
        
        if len(preview_text) > 15:
            content_p.font.size = Pt(9)

def create_assessment_slide(prs, content):
    """Create resources and assessment slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(252, 251, 247)

    header_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.33), Inches(0.6)
    )
    header_fill = header_bar.fill
    header_fill.solid()
    header_fill.fore_color.rgb = RGBColor(139, 134, 128)

    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(11.33), Inches(0.8))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = "STRATEGIC ASSESSMENT"
    title_p.font.size = Pt(28)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(87, 83, 78)
    title_p.alignment = PP_ALIGN.CENTER

    assessments = [
        {
            'title': 'Resources &\nPartnerships',
            'content': content.resources_partnerships or 'Not specified',
            'color': RGBColor(169, 156, 141),  
            'pos': (1, 2.2, 3.7, 4)
        },
        {
            'title': 'Skills &\nCapabilities', 
            'content': content.skills_capabilities or 'Not specified',
            'color': RGBColor(156, 153, 136),  
            'pos': (4.9, 2.2, 3.7, 4)
        },
        {
            'title': 'Impact\nAvenues',
            'content': content.impact_avenues or 'Not specified',
            'color': RGBColor(176, 151, 128), 
            'pos': (8.8, 2.2, 3.7, 4)
        }
    ]
    
    for assessment in assessments:
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(assessment['pos'][0]), Inches(assessment['pos'][1]),
            Inches(assessment['pos'][2]), Inches(assessment['pos'][3])
        )
        
        fill = card.fill
        fill.solid()
        fill.fore_color.rgb = assessment['color']
  
        card.line.color.rgb = RGBColor(255, 254, 251)
        card.line.width = Pt(3)

        text_frame = card.text_frame
        text_frame.margin_left = Inches(0.3)
        text_frame.margin_right = Inches(0.3)
        text_frame.margin_top = Inches(0.3)
        text_frame.margin_bottom = Inches(0.3)

        title_p = text_frame.paragraphs[0]
        title_p.text = assessment['title']
        title_p.font.size = Pt(16)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(255, 254, 251)
        title_p.alignment = PP_ALIGN.CENTER

        content_p = text_frame.add_paragraph()
        content_text = assessment['content'][:150] + "..." if len(assessment['content']) > 150 else assessment['content']
        content_p.text = content_text
        content_p.font.size = Pt(11)
        content_p.font.color.rgb = RGBColor(255, 254, 251)

def create_risk_slide(prs, content):
    """Create risk assessment slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(252, 251, 247)

    header_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.33), Inches(0.6)
    )
    header_fill = header_bar.fill
    header_fill.solid()
    header_fill.fore_color.rgb = RGBColor(139, 134, 128)

    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(11.33), Inches(0.8))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = "RISK ASSESSMENT & MITIGATION"
    title_p.font.size = Pt(28)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(87, 83, 78)
    title_p.alignment = PP_ALIGN.CENTER

    risk_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(2.2), Inches(5.5), Inches(4)
    )
    risk_fill = risk_box.fill
    risk_fill.solid()
    risk_fill.fore_color.rgb = RGBColor(186, 151, 143)  

    risk_box.line.color.rgb = RGBColor(255, 254, 251)
    risk_box.line.width = Pt(3)

    risk_frame = risk_box.text_frame
    risk_frame.margin_left = Inches(0.3)
    risk_frame.margin_top = Inches(0.3)
    risk_frame.margin_right = Inches(0.3)
    risk_frame.margin_bottom = Inches(0.3)
    
    risk_title = risk_frame.paragraphs[0]
    risk_title.text = "RISKS & INHIBITORS"
    risk_title.font.size = Pt(18)
    risk_title.font.bold = True
    risk_title.font.color.rgb = RGBColor(255, 254, 251)
    risk_title.alignment = PP_ALIGN.CENTER
    
    risk_content = risk_frame.add_paragraph()
    risk_text = content.risks_inhibitors or "No risks identified"
    risk_content.text = risk_text[:200] + "..." if len(risk_text) > 200 else risk_text
    risk_content.font.size = Pt(12)
    risk_content.font.color.rgb = RGBColor(255, 254, 251)

    mitigation_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.8), Inches(2.2), Inches(5.5), Inches(4)
    )
    mitigation_fill = mitigation_box.fill
    mitigation_fill.solid()
    mitigation_fill.fore_color.rgb = RGBColor(156, 174, 146)  

    mitigation_box.line.color.rgb = RGBColor(255, 254, 251)
    mitigation_box.line.width = Pt(3)

    mitigation_frame = mitigation_box.text_frame
    mitigation_frame.margin_left = Inches(0.3)
    mitigation_frame.margin_top = Inches(0.3)
    mitigation_frame.margin_right = Inches(0.3)
    mitigation_frame.margin_bottom = Inches(0.3)
    
    mitigation_title = mitigation_frame.paragraphs[0]
    mitigation_title.text = "MITIGATION STRATEGIES"
    mitigation_title.font.size = Pt(18)
    mitigation_title.font.bold = True
    mitigation_title.font.color.rgb = RGBColor(255, 254, 251)
    mitigation_title.alignment = PP_ALIGN.CENTER
    
    mitigation_content = mitigation_frame.add_paragraph()
    mitigation_text = content.mitigation_strategies or "No mitigation strategies defined"
    mitigation_content.text = mitigation_text[:200] + "..." if len(mitigation_text) > 200 else mitigation_text
    mitigation_content.font.size = Pt(12)
    mitigation_content.font.color.rgb = RGBColor(255, 254, 251)

@api_view(['POST'])
@permission_classes([permissions.IsAuthenticated])
def submit_form_for_review(request, form_id):
    """Submit form for review - Only team owners can submit Action/Education forms"""
    form = get_object_or_404(Form, id=form_id)
    
    if form.team is None:
        if form.created_by != request.user:
            return Response(
                {'error': 'Only the form creator can submit personal forms for review'}, 
                status=status.HTTP_403_FORBIDDEN
            )
    else:
        # Check if user is team owner
        membership = TeamMembership.objects.filter(
            user=request.user,
            team=form.team,
            role='owner'
        ).first()
        
        if not membership:
            return Response(
                {'error': 'Only team owner can submit forms for review'}, 
                status=status.HTTP_403_FORBIDDEN
            )
    
    # Check if form type is reviewable
    if form.type not in ['action', 'education']:
        return Response(
            {'error': 'Only Action and Education forms can be submitted for review'}, 
            status=status.HTTP_400_BAD_REQUEST
        )
    
    # Check if form is already submitted or approved
    if form.review_status in ['submitted_for_review', 'under_review', 'approved']:
        return Response(
            {'error': f'Form is already {form.review_status.replace("_", " ")}'}, 
            status=status.HTTP_400_BAD_REQUEST
        )
    
    # Check if form has content
    try:
        content = FormContent.objects.get(form=form)
        if not content.title or not content.description:
            return Response(
                {'error': 'Form must have at least title and description to be submitted'}, 
                status=status.HTTP_400_BAD_REQUEST
            )
    except FormContent.DoesNotExist:
        return Response(
            {'error': 'Form content not found'}, 
            status=status.HTTP_400_BAD_REQUEST
        )
    
    # Update form status
    form.review_status = 'submitted_for_review'
    form.submitted_for_review_at = timezone.now()
    form.save(update_fields=['review_status', 'submitted_for_review_at'])
    
    # Create notifications for all system administrators (users with is_staff=True)
    admin_users = User.objects.filter(
        is_staff=True
    ).exclude(id=request.user.id)
    
    notification_data = {
        'form_id': form.id,
        'form_title': content.title or form.title,
        'form_type': form.type,
        'form_source': 'personal' if form.team is None else 'team',
        'team_id': form.team.id if form.team else None,
        'team_name': form.team.name if form.team else 'Personal Form',
        'submitter_username': request.user.username,
        'submitter_email': request.user.email,
        'submitted_at': form.submitted_for_review_at.isoformat(),
    }
    
    # Create notifications for all administrators
    notifications_created = 0
    for admin_user in admin_users:
        Notification.objects.create(
            recipient=admin_user,
            notification_type='form_review_request',
            data=notification_data,
            expires_at=timezone.now() + timedelta(days=30),
            status='pending'
        )
        notifications_created += 1
    
    return Response({
        'message': 'Form submitted for review successfully',
        'review_status': form.review_status,
        'submitted_at': form.submitted_for_review_at,
        'notifications_sent': notifications_created
    }, status=status.HTTP_200_OK)

@api_view(['GET'])
@permission_classes([permissions.IsAuthenticated])
def get_forms_for_review(request):
    """Get all forms pending review for current user (team owners only) - for team owners to see their own team's forms"""
    # Get teams where user is owner
    owned_teams = TeamMembership.objects.filter(
        user=request.user,
        role='owner'
    ).values_list('team_id', flat=True)
    
    if not owned_teams:
        return Response({
            'forms': [],
            'message': 'You are not a team owner'
        }, status=status.HTTP_200_OK)
    
    # Get forms pending review from owned teams
    pending_forms = Form.objects.filter(
        team_id__in=owned_teams,
        review_status='submitted_for_review',
        type__in=['action', 'education']
    ).select_related('created_by', 'team').prefetch_related('content')
    
    forms_data = []
    for form in pending_forms:
        try:
            content = form.content
            forms_data.append({
                'id': form.id,
                'title': content.title or form.title,
                'description': content.description or form.description,
                'type': form.type,
                'team_name': form.team.name if form.team else 'Personal Form',
                'submitted_by': form.created_by.username,
                'submitted_at': form.submitted_for_review_at,
                'created_at': form.created_at,
                'is_personal': form.team is None,
            })
        except FormContent.DoesNotExist:
            continue
    
    return Response({
        'forms': forms_data,
        'count': len(forms_data)
    }, status=status.HTTP_200_OK)

@api_view(['GET'])
@permission_classes([permissions.IsAuthenticated])
def get_pending_review_forms(request):
    """Get all forms pending review for system administrators (is_staff=True users)"""
    # Check if user is a system administrator (is_staff=True)
    if not request.user.is_staff:
        return Response({
            'forms': [],
            'message': 'You are not a system administrator'
        }, status=status.HTTP_403_FORBIDDEN)
    
    # Get all forms pending review
    pending_forms = Form.objects.filter(
        review_status='submitted_for_review',
        type__in=['action', 'education']
    ).select_related('created_by', 'team').prefetch_related('content')
    
    forms_data = []
    for form in pending_forms:
        try:
            content = form.content
            forms_data.append({
                'id': form.id,
                'title': content.title or form.title,
                'description': content.description or form.description,
                'type': form.type,
                'team_name': form.team.name if form.team else 'Personal Form',
                'team_id': form.team.id if form.team else None,
                'submitted_by': form.created_by.username,
                'submitted_at': form.submitted_for_review_at,
                'created_at': form.created_at,
            })
        except FormContent.DoesNotExist:
            continue
    
    return Response({
        'forms': forms_data,
        'count': len(forms_data)
    }, status=status.HTTP_200_OK)

@api_view(['GET'])
@permission_classes([permissions.IsAuthenticated])
def get_form_for_review_detail(request, form_id):
    """Get detailed form data for review - Only system administrators can review"""
    form = get_object_or_404(Form, id=form_id)
    
    # Check if user is system administrator (is_staff=True)
    if not request.user.is_staff:
        return Response(
            {'error': 'Only system administrators can review forms'}, 
            status=status.HTTP_403_FORBIDDEN
        )
    
    if form.review_status not in ['submitted_for_review', 'under_review']:
        return Response(
            {'error': 'Form is not available for review'}, 
            status=status.HTTP_400_BAD_REQUEST
        )
    
    # Get form content first
    try:
        content = form.content
    except FormContent.DoesNotExist:
        return Response(
            {'error': 'Form content not found'}, 
            status=status.HTTP_404_NOT_FOUND
        )

    # Mark as under review if not already
    if form.review_status == 'submitted_for_review':
        form.review_status = 'under_review'
        form.save(update_fields=['review_status'])

        if form.team is None:
            recipients = [form.created_by]
        else:
            recipients = User.objects.filter(
                team_memberships__team=form.team
            ).exclude(id=request.user.id)
        
        notification_data = {
            'form_id': form.id,
            'form_title': content.title or form.title,
            'form_type': form.type,
            'team_name': form.team.name if form.team else 'Personal Form',
            'reviewer_username': request.user.username,
            'status': 'under_review',
            'message': f'Your form "{content.title or form.title}" is now under review.'
        }
        
        for recipient in recipients:
            Notification.objects.create(
                recipient=recipient,
                notification_type='form_review_status_update',
                data=notification_data,
                expires_at=timezone.now() + timedelta(days=7),
                status='pending'
            )
    
    serializer = FormContentSerializer(content)
    
    return Response({
        'form': {
            'id': form.id,
            'title': form.title,
            'type': form.type,
            'team_name': form.team.name if form.team else 'Personal Form',
            'submitted_by': form.created_by.username,
            'submitted_at': form.submitted_for_review_at,
            'review_status': form.review_status,
        },
        'content': serializer.data
    }, status=status.HTTP_200_OK)
        

@api_view(['POST'])
@permission_classes([permissions.IsAuthenticated])
def approve_form_review(request, form_id):
    """Approve form and transfer to main database - Only system administrators can approve"""
    form = get_object_or_404(Form, id=form_id)
    
    # Check if user is system administrator (is_staff=True)
    if not request.user.is_staff:
        return Response(
            {'error': 'Only system administrators can approve forms'}, 
            status=status.HTTP_403_FORBIDDEN
        )
    
    if form.review_status not in ['submitted_for_review', 'under_review']:
        return Response(
            {'error': 'Form is not available for approval'}, 
            status=status.HTTP_400_BAD_REQUEST
        )
    
    try:
        content = form.content
        comments = request.data.get('comments', '')
        
        # Transfer to main database based on form type
        if form.type == 'action':
            action_record = ActionDb.objects.create(
                actions=content.actions or '',
                action_detail=content.action_detail or '',
                field_sdgs=content.sdgs_related or '',
                level=str(content.level) if content.level is not None else '',
                individual_organization=content.individual_organization,
                location_specific_actions_org_onlyonly_field=content.location or '',
                related_industry_org_only_field=content.related_industry or '',
                digital_actions=1 if content.digital_actions else 0,
                source_descriptions=content.source_descriptions or '',
                award=1 if content.award else 0,
                source_links=content.source_links or '',
                additional_notes=content.additional_notes or '',
                award_descriptions=content.award_descriptions or '',
            )
            record_id = action_record.id
            
        elif form.type == 'education':
            education_record = EducationDb.objects.create(
                title=content.title or '',
                descriptions=content.description or '',
                aims=content.aims or '',
                learning_outcome_expecting_outcome_field=content.learning_outcomes or '',
                type_label=content.type_label or '',
                location=content.location or '',
                organization=content.organization or '',
                year=content.year or '',
                sdgs_related=content.sdgs_related or '',
                related_to_which_discipline=content.related_discipline or '',
                useful_for_which_industries=content.useful_industries or '',
                source=content.source or '',
                link=content.link or '',
            )
            record_id = education_record.id
        
        # Update form status
        form.review_status = 'approved'
        form.reviewed_by = request.user
        form.reviewed_at = timezone.now()
        form.review_comments = comments
        form.save(update_fields=['review_status', 'reviewed_by', 'reviewed_at', 'review_comments'])
        
        Notification.objects.filter(
            notification_type='form_review_request',
            data__form_id=form.id,
            status='pending'
        ).update(status='accepted')
        
        if form.team is None:
            recipients = [form.created_by]
        else:
            recipients = User.objects.filter(
                team_memberships__team=form.team
            ).exclude(id=request.user.id)
        
        for recipient in recipients:
            Notification.objects.create(
                recipient=recipient,
                notification_type='form_review_completed',
                data={
                    'form_id': form.id,
                    'form_title': content.title or form.title,
                    'form_type': form.type,
                    'team_name': form.team.name if form.team else 'Personal Form',
                    'reviewer_username': request.user.username,
                    'reviewed_at': form.reviewed_at.isoformat(),
                    'status': 'approved',
                    'comments': comments,
                    'main_db_id': record_id,
                    'message': f'Great news! Your form "{content.title or form.title}" has been approved and added to the main database.'
                },
                expires_at=timezone.now() + timedelta(days=30),
                status='pending'
            )
        
        return Response({
            'message': 'Form approved and added to main database successfully',
            'review_status': form.review_status,
            'main_db_id': record_id,
            'reviewed_at': form.reviewed_at
        }, status=status.HTTP_200_OK)
        
    except FormContent.DoesNotExist:
        return Response(
            {'error': 'Form content not found'}, 
            status=status.HTTP_404_NOT_FOUND
        )
    except Exception as e:
        return Response(
            {'error': f'Failed to approve form: {str(e)}'}, 
            status=status.HTTP_500_INTERNAL_SERVER_ERROR
        )

@api_view(['POST'])
@permission_classes([permissions.IsAuthenticated])
def reject_form_review(request, form_id):
    """Reject form review - Only system administrators can reject"""
    form = get_object_or_404(Form, id=form_id)
    
    # Check if user is system administrator (is_staff=True)
    if not request.user.is_staff:
        return Response(
            {'error': 'Only system administrators can reject forms'}, 
            status=status.HTTP_403_FORBIDDEN
        )
    
    if form.review_status not in ['submitted_for_review', 'under_review']:
        return Response(
            {'error': 'Form is not available for rejection'}, 
            status=status.HTTP_400_BAD_REQUEST
        )
    
    comments = request.data.get('comments', '')
    if not comments:
        return Response(
            {'error': 'Rejection reason is required'}, 
            status=status.HTTP_400_BAD_REQUEST
        )
    
    try:
        content = form.content
        
        # Update form status
        form.review_status = 'rejected'
        form.reviewed_by = request.user
        form.reviewed_at = timezone.now()
        form.review_comments = comments
        form.save(update_fields=['review_status', 'reviewed_by', 'reviewed_at', 'review_comments'])
        
        Notification.objects.filter(
            notification_type='form_review_request',
            data__form_id=form.id,
            status='pending'
        ).update(status='rejected')
        
        if form.team is None:
            recipients = [form.created_by]
        else:
            recipients = User.objects.filter(
                team_memberships__team=form.team
            ).exclude(id=request.user.id)
        
        for recipient in recipients:
            Notification.objects.create(
                recipient=recipient,
                notification_type='form_review_completed',
                data={
                    'form_id': form.id,
                    'form_title': content.title or form.title,
                    'form_type': form.type,
                    'team_name': form.team.name if form.team else 'Personal Form',
                    'reviewer_username': request.user.username,
                    'reviewed_at': form.reviewed_at.isoformat(),
                    'status': 'rejected',
                    'comments': comments,
                    'message': f'Your form "{content.title or form.title}" needs revision. Please check the feedback and resubmit.'
                },
                expires_at=timezone.now() + timedelta(days=30),
                status='pending'
            )
        
        return Response({
            'message': 'Form rejected successfully',
            'review_status': form.review_status,
            'reviewed_at': form.reviewed_at,
            'comments': comments
        }, status=status.HTTP_200_OK)
        
    except FormContent.DoesNotExist:
        return Response(
            {'error': 'Form content not found'}, 
            status=status.HTTP_404_NOT_FOUND
        )
    
class PersonalFormListCreateView(generics.ListCreateAPIView):
    """List and create personal forms for the authenticated user"""
    permission_classes = [permissions.IsAuthenticated]
    
    def get_queryset(self):
        """Get personal forms created by the authenticated user"""
        return Form.objects.filter(
            team__isnull=True,
            created_by=self.request.user
        ).select_related('created_by', 'last_modified_by')
    
    def get_serializer_class(self):
        if self.request.method == 'POST':
            return CreatePersonalFormSerializer
        return FormListSerializer
    
    def create(self, request, *args, **kwargs):
        """Create a new personal form"""
        return super().create(request, *args, **kwargs)

@api_view(['GET'])
@permission_classes([permissions.IsAuthenticated])
def personal_form_stats(request):
    """Get statistics for personal forms created by the authenticated user"""
    forms = Form.objects.filter(
        team__isnull=True,
        created_by=request.user
    )
    
    stats = {
        'total_forms': forms.count(),
        'active_forms': forms.filter(status='active').count(),
        'locked_forms': forms.filter(status='locked').count(),
        'archived_forms': forms.filter(status='archived').count(),
        'forms_by_type': {
            'action': forms.filter(type='action').count(),
            'education': forms.filter(type='education').count(),
            'blank': forms.filter(type='blank').count(),
            'ida': forms.filter(type='ida').count(),
        }
    }
    
    return Response(stats)